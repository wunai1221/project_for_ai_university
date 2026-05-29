"""
PPT 抽取程式
功能：抽取每張投影片的文字 + OCR 辨識圖片文字
輸出位置：rag/data/raw/pptx/
"""

import json
import os
import io
import numpy as np
from datetime import datetime
from pptx import Presentation
from PIL import Image
import easyocr

# ============================================================
# 設定區
# ============================================================
PPTX_FILES = [
    {
        "filename": "鉅昕集團介紹(完整).pptx",
        "category_hint": ["鋼筋加工", "鋼構工程", "鋼材買賣及加工", "太陽能統包設計到送電完成", "餐廳推廣", "農產品推銷"]
    },
    {
        "filename": "最適合安裝廠房屋頂的不銹鋼鋼瓦太陽能系統.pptx",
        "category_hint": ["太陽能統包設計到送電完成"]
    },
]

INPUT_DIR = os.path.join(os.path.dirname(__file__), "../../data/raw/pptx")
OUTPUT_DIR = os.path.join(os.path.dirname(__file__), "../../data/raw/pptx")
MIN_CHARS = 10
# ============================================================

# 初始化 OCR（只跑一次）
print("載入 OCR 模型中...")
reader = easyocr.Reader(["ch_tra", "en"])
print("OCR 模型載入完成！")


def extract_text_from_slide(slide) -> str:
    texts = []
    for shape in slide.shapes:
        # 一般文字
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                line = para.text.strip()
                if line:
                    texts.append(line)
        # 表格
        if shape.has_table:
            for row in shape.table.rows:
                row_texts = []
                for cell in row.cells:
                    cell_text = cell.text.strip()
                    if cell_text:
                        row_texts.append(cell_text)
                if row_texts:
                    texts.append(" | ".join(row_texts))
    return "\n".join(texts)


def extract_images_from_slide(slide) -> list:
    """把投影片裡的圖片抽出來，回傳 PIL Image list"""
    images = []
    for shape in slide.shapes:
        if shape.shape_type == 13:
            image_data = shape.image.blob
            image = Image.open(io.BytesIO(image_data)).convert("RGB")
            images.append(image)
    return images


def ocr_image(image: Image.Image) -> str:
    """對一張圖片跑 OCR，回傳辨識出的文字"""
    img_array = np.array(image)
    result = reader.readtext(img_array)
    if not result:
        return ""
    lines = [item[1] for item in result if item[2] > 0.5]
    return "\n".join(lines)


def extract_pptx(filepath: str, file_info: dict) -> list[dict]:
    """處理整份 PPT，每張投影片輸出一筆資料"""
    prs = Presentation(filepath)
    filename = file_info["filename"]
    category_hint = file_info["category_hint"]
    results = []

    for i, slide in enumerate(prs.slides):
        slide_num = i + 1
        print(f"  處理第 {slide_num} 張投影片...")

        text = extract_text_from_slide(slide)

        images = extract_images_from_slide(slide)
        ocr_texts = []
        for img in images:
            ocr_result = ocr_image(img)
            if ocr_result:
                ocr_texts.append(ocr_result)

        combined = text
        if ocr_texts:
            combined += "\n[圖片文字]\n" + "\n".join(ocr_texts)

        if len(combined.strip()) < MIN_CHARS:
            print(f"    → 內容太少，跳過")
            continue

        results.append({
            "text": combined.strip(),
            "source_file": filename,
            "source_type": "pptx",
            "page_or_section": f"第{slide_num}張投影片",
            "category_hint": category_hint,
            "category": "",
            "date_processed": datetime.now().strftime("%Y-%m-%d")
        })

    return results


def main():
    all_results = []

    for file_info in PPTX_FILES:
        filepath = os.path.join(INPUT_DIR, file_info["filename"])

        if not os.path.exists(filepath):
            print(f"找不到檔案：{filepath}，請確認檔案是否放在 data/raw/pptx/ 資料夾")
            continue

        print(f"\n處理：{file_info['filename']}")
        results = extract_pptx(filepath, file_info)
        print(f"  → 共處理 {len(results)} 張投影片")

        out_name = file_info["filename"].replace(".pptx", ".json")
        out_path = os.path.join(OUTPUT_DIR, out_name)
        with open(out_path, "w", encoding="utf-8") as f:
            json.dump(results, f, ensure_ascii=False, indent=2)

        all_results.extend(results)

    merged_path = os.path.join(OUTPUT_DIR, "_all_pptx.json")
    with open(merged_path, "w", encoding="utf-8") as f:
        json.dump(all_results, f, ensure_ascii=False, indent=2)

    print(f"\n全部完成！共處理 {len(all_results)} 張投影片")


if __name__ == "__main__":
    main()